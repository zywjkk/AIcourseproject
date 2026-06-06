from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[4]
OUT = ROOT / "outputs" / "manual-ppt-tomato" / "presentations" / "tomato-defense" / "output"
FINAL = OUT / "test开题报告test番茄病害识别答辩PPT.pptx"
EXP = ROOT / "outputs" / "experiments"
BEST = EXP / "mobilenet_v3_transfer_fc_then_features_tail"


class Theme:
    bg = RGBColor(250, 250, 247)
    ink = RGBColor(23, 32, 26)
    muted = RGBColor(100, 112, 103)
    faint = RGBColor(220, 226, 218)
    pale = RGBColor(238, 243, 236)
    green = RGBColor(47, 107, 79)
    green2 = RGBColor(95, 141, 106)
    gold = RGBColor(181, 138, 37)
    red = RGBColor(164, 78, 66)
    white = RGBColor(255, 255, 255)


W, H = Inches(13.333), Inches(7.5)
FONT = "Microsoft YaHei"
FONT_EN = "Aptos"


def rgb_to_hex(rgb):
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def add_text(slide, text, x, y, w, h, size=14, color=Theme.ink, bold=False,
             align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=Theme.white, line=Theme.faint, width=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(width)
    return shp


def base(slide, n, section="番茄病害识别检测系统"):
    add_rect(slide, 0, 0, 13.333, 7.5, Theme.bg, Theme.bg, 0)
    add_rect(slide, 0.5, 0.52, 12.35, 0.01, Theme.faint, Theme.faint, 0)
    add_text(slide, section, 0.52, 0.25, 6.2, 0.22, 8.5, Theme.muted)
    add_text(slide, f"{n:02d}", 12.15, 0.22, 0.7, 0.26, 10, Theme.green, True, PP_ALIGN.RIGHT, FONT_EN)


def slide_title(slide, title, subtitle=""):
    add_text(slide, title, 0.65, 0.75, 7.8, 0.5, 24, Theme.ink, True)
    if subtitle:
        add_text(slide, subtitle, 0.67, 1.28, 9.4, 0.32, 10.5, Theme.muted)


def metric(slide, value, label, x, y, w=1.75):
    add_rect(slide, x, y, w, 0.82, Theme.white, Theme.faint, 0.75)
    add_text(slide, value, x + 0.12, y + 0.11, w - 0.24, 0.32, 20, Theme.green, True, font=FONT_EN)
    add_text(slide, label, x + 0.12, y + 0.51, w - 0.24, 0.20, 7.8, Theme.muted)


def bullets(slide, items, x, y, w, gap=0.34, size=11, color=Theme.ink):
    for i, item in enumerate(items):
        text = item if isinstance(item, str) else item[0]
        dot = Theme.green if isinstance(item, str) or len(item) < 2 else item[1]
        yy = y + i * gap
        add_rect(slide, x, yy + 0.09, 0.055, 0.055, dot, dot, 0)
        add_text(slide, text, x + 0.18, yy, w - 0.18, gap * 0.95, size, color)


def table(slide, headers, rows, x, y, widths, row_h=0.38):
    total = sum(widths)
    add_rect(slide, x, y, total, row_h, Theme.green, Theme.green, 0)
    cx = x
    for i, head in enumerate(headers):
        add_text(slide, head, cx + 0.06, y + 0.07, widths[i] - 0.12, row_h - 0.08, 8.2, Theme.white, True)
        cx += widths[i]
    for r, row in enumerate(rows):
        yy = y + row_h * (r + 1)
        add_rect(slide, x, yy, total, row_h, Theme.white if r % 2 == 0 else Theme.pale, Theme.faint, 0.35)
        cx = x
        for c, cell in enumerate(row):
            add_text(slide, str(cell), cx + 0.06, yy + 0.07, widths[c] - 0.12, row_h - 0.08,
                     7.9, Theme.ink if c == 0 else Theme.muted, c == 0)
            cx += widths[c]


def fit_picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def bar(slide, label, value, x, y, w, color=Theme.green):
    add_text(slide, label, x, y - 0.04, 1.65, 0.18, 7.7, Theme.ink)
    add_rect(slide, x + 1.72, y, w, 0.09, Theme.faint, Theme.faint, 0)
    add_rect(slide, x + 1.72, y, w * value, 0.09, color, color, 0)
    add_text(slide, f"{value * 100:.1f}%", x + 1.78 + w, y - 0.05, 0.55, 0.18, 7.5, color, True, font=FONT_EN)


def node(slide, text, x, y, w=1.55, h=0.55, fill=Theme.white):
    add_rect(slide, x, y, w, h, fill, Theme.green2, 0.75)
    add_text(slide, text, x + 0.06, y + 0.08, w - 0.12, h - 0.12, 8.0, Theme.ink, False, PP_ALIGN.CENTER)


def build():
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    base(s, 1, "人工智能基础(A) / 选题四")
    add_rect(s, 0.55, 0.9, 0.07, 1.55, Theme.green, Theme.green, 0)
    add_text(s, "番茄病害识别检测系统", 0.82, 0.86, 8.9, 0.55, 27, Theme.ink, True)
    add_text(s, "深度学习在农业中的应用 / CNN + Transfer Learning + Gradio Demo", 0.84, 1.5, 8.8, 0.28, 11.5, Theme.muted, font=FONT_EN)
    add_text(s, "PlantVillage 番茄叶片 10 类分类任务；完整覆盖数据预处理、模型训练、迁移学习、评估可视化与交互部署。",
             0.84, 2.12, 7.3, 0.55, 14, Theme.ink)
    for args in [("96.6%", "测试集 Top-1 Accuracy", 0.84, 3.45, 1.9), ("10", "番茄病害/健康类别", 3.05, 3.45, 1.9),
                 ("11k", "PlantVillage 图像规模", 5.25, 3.45, 1.9), ("Top-3", "网页实时预测输出", 7.45, 3.45, 1.9)]:
        metric(s, *args)
    add_rect(s, 0.84, 5.82, 10.7, 0.01, Theme.faint, Theme.faint, 0)
    add_text(s, "答辩重点：核心原理与算法实现；不展开宏观背景，直接展示工程流程和实验结果。",
             0.84, 6.08, 9.4, 0.26, 10.5, Theme.muted)

    s = prs.slides.add_slide(blank)
    base(s, 2)
    slide_title(s, "课程要求对齐", "把评分点转化为可检查的代码、实验和演示证据")
    table(s, ["课程要求", "项目实现", "证据文件/结果"],
          [["完整训练脚本", "Dataset/DataLoader、模型、训练循环、验证评估", "train.py + src/*"],
           ["Top-1 ≥ 90%", "最佳 MobileNetV3-small 迁移学习", "96.6%"],
           ["CNN 与 Shape", "从输入到分类头的 Tensor Shape 记录", "model_architecture.txt"],
           ["迁移学习冻结", "fc_only / fc_then_layer4 / scratch 对比", "6 组实验"],
           ["评估可视化", "Loss/Acc、混淆矩阵、分类报告", "outputs/experiments/*"],
           ["进阶挑战", "Gradio Top-3 预测；3 骨干对比", "app.py + 对比图"]],
          0.72, 1.82, [2.35, 4.35, 3.45], 0.45)
    add_text(s, "交付逻辑：源码负责可复现，实验报告负责可解释，PPT/演示视频负责可展示。",
             0.72, 5.65, 9.6, 0.32, 14, Theme.green, True)

    s = prs.slides.add_slide(blank)
    base(s, 3)
    slide_title(s, "数据集与预处理", "PlantVillage Tomato：10 类病害/健康叶片，固定随机种子确保复现")
    for args in [("10000", "Train", 0.72, 1.95, 1.55), ("500", "Validation", 2.5, 1.95, 1.55), ("500", "Test", 4.28, 1.95, 1.55)]:
        metric(s, *args)
    add_text(s, "划分策略", 0.72, 3.08, 1.4, 0.22, 10.5, Theme.green, True)
    bullets(s, ["原始 train 作为训练集；原始 val 按固定 seed 1:1 拆分为验证集与测试集",
                "固定 Python / NumPy / PyTorch 随机种子，报告中同步说明复现设置",
                "类别顺序由 ImageFolder 读取并保存，避免预测标签与训练标签错位"],
            0.72, 3.48, 5.7, 0.38, 10.7)
    add_rect(s, 7.2, 1.9, 4.35, 3.28, Theme.white, Theme.faint, 0.75)
    add_text(s, "Transform Pipeline", 7.45, 2.15, 2.2, 0.25, 12, Theme.ink, True, font=FONT_EN)
    bullets(s, ["Resize: 224 × 224", "RandomHorizontalFlip", "RandomRotation: ±15°", "ToTensor + ImageNet Normalize", "Valid/Test 仅确定性预处理"],
            7.45, 2.65, 3.5, 0.43, 11.4)

    s = prs.slides.add_slide(blank)
    base(s, 4)
    slide_title(s, "CNN 核心结构与 Tensor Shape", "从局部纹理到病害类别：卷积特征逐层抽象，分类头输出 10 维概率")
    nodes = [("Input\n[B,3,224,224]", 0.75, 2.28, 1.45), ("Conv/BN/ReLU\n[B,64,112,112]", 2.55, 2.28, 1.75),
             ("MaxPool+Blocks\n[B,64,56,56]", 4.65, 2.28, 1.75), ("Layer2-4\n[B,512,7,7]", 6.75, 2.28, 1.55),
             ("AvgPool+Linear\n[B,10]", 8.6, 2.28, 1.65)]
    for i, (t, x, y, w) in enumerate(nodes):
        node(s, t, x, y, w, 0.78, Theme.pale if i == len(nodes) - 1 else Theme.white)
        if i < len(nodes) - 1:
            add_rect(s, x + w + 0.15, y + 0.38, 0.32, 0.02, Theme.green2, Theme.green2, 0)
    add_text(s, "关键组件", 0.75, 3.78, 1.4, 0.22, 10.5, Theme.green, True)
    add_text(s, "Conv2D 提取叶片斑点、边缘和颜色纹理；MaxPool 降低空间分辨率并增强平移鲁棒性；ReLU 引入非线性；Dropout/权重衰减缓解过拟合；Linear 将高层语义映射到 10 类。",
             0.75, 4.16, 10.4, 0.7, 13, Theme.ink)
    add_text(s, "报告中保存模型结构、参数量和 Tensor Shape 变换，对应课程“模型架构图/Shape”要求。",
             0.75, 5.55, 9.5, 0.3, 11.5, Theme.muted)

    s = prs.slides.add_slide(blank)
    base(s, 5)
    slide_title(s, "迁移学习与冻结策略", "冻结不是装饰项：它决定哪些 ImageNet 特征被保留，哪些层适配番茄病害")
    table(s, ["策略", "训练参数", "目的", "测试准确率"],
          [["From Scratch", "全模型", "观察无预训练基线", "93.0%"],
           ["fc_only", "仅最后分类层", "验证“冻结全部 backbone”", "89.0%"],
           ["fc_then_layer4", "分类层 + 尾部特征层", "冻结前面 N 层，微调高层语义", "95.4%~96.6%"]],
          0.72, 1.88, [2.1, 2.35, 3.95, 1.3], 0.55)
    add_rect(s, 1.05, 4.45, 8.9, 0.78, Theme.white, Theme.faint, 0.75)
    add_text(s, "冻结前面 N 层", 1.28, 4.72, 1.4, 0.22, 10.2, Theme.green, True)
    for x, txt, fill, col in [(3.0, "浅层纹理", Theme.faint, Theme.muted), (4.65, "中层形状", Theme.faint, Theme.muted),
                              (6.3, "高层语义", Theme.pale, Theme.green), (7.95, "分类头", Theme.green, Theme.white)]:
        add_rect(s, x, 4.68, 1.25, 0.26, fill, fill, 0)
        add_text(s, txt, x, 4.72, 1.25, 0.16, 7.2, col, False, PP_ALIGN.CENTER)
    add_text(s, "结论：仅训练分类头不足以适配病害细粒度差异；解冻尾部特征层后，迁移学习明显优于从头训练和 fc_only。",
             1.05, 5.6, 9.1, 0.35, 12.7, Theme.ink, True)

    s = prs.slides.add_slide(blank)
    base(s, 6)
    slide_title(s, "工程实现结构", "按 PyTorch 项目组织拆分，训练、评估、部署互相解耦")
    module_nodes = [("src/config.py\n实验配置 / seed", 0.95, 1.95), ("src/data.py\nDataset / DataLoader", 3.3, 1.95),
                    ("src/models.py\nBackbone / Freeze", 5.65, 1.95), ("src/trainer.py\n训练循环 / 验证", 8.0, 1.95),
                    ("src/evaluator.py\n报告 / 混淆矩阵", 2.1, 3.35), ("src/gradcam.py\n可解释性热力图", 4.45, 3.35),
                    ("train.py / evaluate.py\n命令行入口", 6.8, 3.35), ("app.py\nGradio Top-3", 9.15, 3.35)]
    for t, x, y in module_nodes:
        node(s, t, x, y, 1.95, 0.66)
    add_text(s, "核心工程原则", 0.95, 5.1, 1.5, 0.22, 10.5, Theme.green, True)
    bullets(s, ["训练脚本可复现：固定随机种子、保存 config/history/metrics/model",
                "评估脚本独立运行：测试集分类报告、混淆矩阵、推理速度统计",
                "Web 演示复用同一 checkpoint 与 transform，避免训练/部署预处理不一致"],
            0.95, 5.48, 9.5, 0.34, 10.8)

    s = prs.slides.add_slide(blank)
    base(s, 7)
    slide_title(s, "实验矩阵与关键结果", "覆盖 Scratch、冻结全部、冻结前 N 层、学习率和 3 种骨干网络对比")
    table(s, ["实验", "骨干", "冻结策略", "Acc", "参数量"],
          [["MobileNetV3", "small", "fc_then_layer4", "96.6%", "1.53M"],
           ["ResNet50", "resnet50", "fc_then_layer4", "96.2%", "23.53M"],
           ["ResNet18 lr=3e-4", "resnet18", "fc_then_layer4", "95.8%", "11.18M"],
           ["ResNet18 lr=1e-3", "resnet18", "fc_then_layer4", "95.4%", "11.18M"],
           ["ResNet18 Scratch", "resnet18", "none", "93.0%", "11.18M"],
           ["ResNet18 fc_only", "resnet18", "fc_only", "89.0%", "11.18M"]],
          0.62, 1.72, [2.35, 1.45, 2.25, 0.9, 1.1], 0.36)
    add_text(s, "Accuracy Ranking", 8.6, 1.72, 1.9, 0.24, 10.5, Theme.green, True, font=FONT_EN)
    for i, (lab, val, col) in enumerate([("MobileNetV3", .966, Theme.green), ("ResNet50", .962, Theme.green2),
                                         ("ResNet18 low LR", .958, Theme.green2), ("ResNet18", .954, Theme.green2),
                                         ("Scratch", .93, Theme.gold), ("fc_only", .89, Theme.red)]):
        bar(s, lab, val, 8.55, 2.15 + i * 0.39, 2.3, col)
    add_text(s, "最终选择 MobileNetV3-small：准确率最高、参数量最低、单张推理最快，适合课程项目中的轻量部署演示。",
             0.62, 5.75, 10.7, 0.34, 13.2, Theme.ink, True)

    s = prs.slides.add_slide(blank)
    base(s, 8)
    slide_title(s, "训练过程与测试集诊断", "Loss/Accuracy 曲线、混淆矩阵和分类报告共同证明模型效果")
    fit_picture(s, BEST / "training_curves.png", 0.62, 1.7, 5.2, 3.25)
    fit_picture(s, BEST / "confusion_matrix.png", 6.35, 1.7, 4.35, 3.25)
    add_rect(s, 10.88, 1.78, 1.45, 3.05, Theme.white, Theme.faint, 0.75)
    add_text(s, "Best model", 11.05, 2.0, 1.05, 0.22, 10, Theme.green, True, font=FONT_EN)
    bullets(s, [("Top-1 96.6%", Theme.green), ("Macro-F1 96.5%", Theme.green), ("healthy F1=1.00", Theme.green),
                ("Septoria F1=0.93", Theme.gold), ("500 张测试图", Theme.muted)], 11.05, 2.45, 1.1, 0.42, 8.5)
    add_text(s, "测试集中主要误差集中在病斑纹理相近的类别；整体 Precision / Recall / F1 均达到约 0.97。",
             0.62, 5.65, 10.1, 0.3, 12.5, Theme.ink)

    s = prs.slides.add_slide(blank)
    base(s, 9)
    slide_title(s, "可解释性与交互演示", "Grad-CAM 验证模型关注叶片病害区域；Gradio 支持上传图片实时 Top-3 预测")
    fit_picture(s, BEST / "grad_cam.png", 0.72, 1.76, 4.55, 3.35)
    add_rect(s, 5.9, 1.85, 5.45, 2.95, Theme.white, Theme.faint, 0.75)
    add_text(s, "Gradio Demo 验证", 6.18, 2.15, 2.4, 0.28, 14, Theme.green, True)
    bullets(s, ["上传番茄叶片图片，返回 Top-3 类别和置信度",
                "Healthy 样本预测：Tomato___healthy，置信度 0.999996",
                "Yellow Leaf Curl Virus 样本预测正确，置信度 0.999882",
                "部署脚本 app.py 复用训练时的类别映射与图像预处理"],
            6.18, 2.72, 4.45, 0.45, 11.4)
    add_text(s, "意义：不仅能给出分类结果，还能解释模型依据，便于答辩中说明模型不是随机猜测。",
             0.72, 5.75, 10.4, 0.32, 12.6, Theme.ink, True)

    s = prs.slides.add_slide(blank)
    base(s, 10)
    slide_title(s, "结论与答辩要点", "项目严格覆盖课程要求，并完成进阶挑战中的交互界面与多骨干对比")
    for args in [("达成", "Top-1 ≥ 90%", 0.82, 1.9, 1.7), ("完整", "训练-评估-部署流程", 2.85, 1.9, 2.0),
                 ("清晰", "模块化 PyTorch 结构", 5.2, 1.9, 2.0), ("加分", "Gradio + 3 Backbone", 7.55, 1.9, 2.1)]:
        metric(s, *args)
    add_text(s, "5分钟展示顺序", 0.88, 3.35, 1.6, 0.24, 11, Theme.green, True)
    bullets(s, ["先说明任务与要求对齐：10 类、90% 指标、不可 AutoML",
                "再讲核心方法：CNN Shape 流转、迁移学习冻结策略",
                "随后展示工程：模块拆分、训练脚本、评估脚本、Web Demo",
                "最后用结果收束：96.6% Accuracy、Macro-F1 96.5%、推理速度和 Grad-CAM"],
            0.88, 3.85, 8.3, 0.43, 12.0)
    add_rect(s, 0.88, 6.02, 10.1, 0.01, Theme.faint, Theme.faint, 0)
    add_text(s, "一句话总结：该系统完成了番茄叶片病害识别从数据、模型、训练、评估到交互部署的完整计算机视觉工程闭环。",
             0.88, 6.28, 9.8, 0.35, 14, Theme.ink, True)

    prs.save(FINAL)
    print(FINAL)


if __name__ == "__main__":
    build()
